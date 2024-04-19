import csv
import io
import os
import uuid
from http.client import HTTPException
import random
import json
import PyPDF2
import docx
import fastapi
from openai import OpenAI
import pandas as pd
from pptx import Presentation
import base64
import tiktoken

from iris.src.utilities.rag_utilities import embedd, load_pickle, save_pickle

accepted_files = {
    "txt": "text",
    "pdf": "PDF",
    "doc": "DOC",
    "csv": "Comma Seperated Values",
    "xls": "Excel Workbooks",
    "ppt": "Powerpoint Presentations"
}


def voice_transcription(audio_file, api_key):
    client = OpenAI(api_key=api_key)
    try:
        audio = open(audio_file, "rb")
        transcript = client.audio.transcriptions.create(
            model="whisper-1",
            file=audio,
            response_format="text",
            language="en"
        )
        return transcript
    except Exception as exc:
        print(exc)
        return exc


def text_to_speech(text, api_key):
    client = OpenAI(api_key=api_key)
    if len(text) > 4096:
        text = text[0:4096]
    response = client.audio.speech.create(
        model="tts-1",
        voice="alloy",
        input=text,
        response_format="wav"
    )
    random_number = random.randrange(1, 10000, 3)
    response.stream_to_file(f"iris/static/audio/{random_number}.wav")
    return f"audio/{random_number}.wav"


def image_generation(prompt, api_key):
    client = OpenAI(api_key=api_key)
    response = client.images.generate(
        model="dall-e-3",
        prompt=prompt,
        size="1024x1024",
        quality="standard",
        n=1,
    )
    image_url = response.data[0].url
    return image_url


def get_audio_file(filename):
    return fastapi.responses.FileResponse(f"static/audio/{filename}.mp3")


def customized_response(prompt, history_log, api_key, temp=0.05, max_tokens=4000, freq_pen=0.0, presc_pen=0.0,
                        url="tbd"):
    history = json.loads(history_log)
    new_prompt = {"role": "user", "content": prompt}
    history.append(new_prompt)
    client = OpenAI(api_key=api_key)
    response = client.chat.completions.create(
        model="gpt-4-0125-preview",
        messages=history
    )
    content = response.choices[0].message.content
    return content


def image_to_text(prompt, image, api_key):
    base64_image = encode_image(image)
    client = OpenAI(api_key=api_key)
    response = client.chat.completions.create(
        model="gpt-4-vision-preview",
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{base64_image}",
                        },
                    },
                ],
            }
        ],
        max_tokens=2000,
    )
    return response.choices[0]


def encode_image(image):
    return base64.b64encode(image.file.read()).decode('utf-8')


async def file_checker(file, api_key="", rag="n"):
    file_content = None
    extension = file.filename.split(".")[-1].lower()
    for file_type in accepted_files.keys():
        if file_type in extension:
            file_content = await file_reader(file, api_key, rag)
            break
    if not file_content:
        raise HTTPException(
            status_code=400,
            detail="File Not supported."
        )
    return file_content


async def file_reader(file, api_key, rag):
    extension = file.filename.split(".")[-1]
    content = await file.read()
    if "pdf" in extension.lower():
        return custom_pdf_reader(file.file, api_key, rag)
    if "doc" in extension.lower():
        return custom_ms_word_reader(file, extension, content, api_key, rag)
    if "csv" in extension.lower():
        return custom_csv_reader(content, api_key, rag)
    if "xls" in extension.lower() or "xlsx" in extension.lower():
        return custom_excel_reader(content, api_key, rag)
    if "ppt" in extension.lower():
        return custom_ppt_reader(file, extension, content, api_key, rag)
    if "txt" in extension.lower():
        if rag == "y":
            embedd_each(content, api_key)
        return content


def embedd_each(file_content, api_key):
    try:
        title = file_content[:50]
        print(title)
        embedded_title = embedd(title)
        embedded_content = embedd(file_content.replace("\n", ""))
        print(embedded_content)
        unique_id = str(uuid.uuid4())
        files = load_pickle()
        files.update({unique_id: {
            "title": title,
            "content": file_content.replace("\n", ""),
            "content_embeddings": embedded_content,
            "title_embeddings": embedded_title
        }})
        save_pickle(files)
    except Exception as exc:
        print("Failed to embed portion of content")


def custom_pdf_reader(pdf_file, api_key, rag):
    pdf = PyPDF2.PdfReader(pdf_file)
    num_pages = len(pdf.pages)
    text = ""
    for i in range(num_pages):
        if rag == "y":
            embedd_each(pdf.pages[i].extract_text(), api_key)
        else:
            text = '\n'.join([text, pdf.pages[i].extract_text()])
    return text


def custom_ms_word_reader(word_file, word_file_extension, word_file_content, api_key, rag):
    doc_id = hash(word_file.filename)
    with open(f"{doc_id}.{word_file_extension}", "wb") as outFile:
        outFile.write(word_file_content)
    doc = docx.Document(f"{doc_id}.{word_file_extension}")
    num_paras = len(doc.paragraphs)
    text = ""
    for i in range(num_paras):
        text = '\n'.join([text, doc.paragraphs[i].text])
    os.remove(f"{doc_id}.{word_file_extension}")
    if rag == "y":
        content_split = text.split("\n\n\n\n")
        if len(content_split) > 1:
            for split in content_split:
                embedd_each(split, api_key)
        else:
            embedd_each(text, api_key)
    return text


def custom_csv_reader(csv_content, api_key, rag):
    data = pd.read_csv(io.BytesIO(csv_content))
    if rag == "y":
        for index, row in data.iterrows():
            row_text = ""
            for column in row:
                row_text = row_text + str(column) + " "
            embedd_each(row_text, api_key)
    return list(csv.reader(csv_content.decode('utf-8').splitlines(), delimiter=','))


def custom_excel_reader(excel_content, api_key, rag):
    data = pd.read_excel(io.BytesIO(excel_content))
    if rag == "y":
        for index, row in data.iterrows():
            row_text = ""
            for column in row:
                row_text = row_text + str(column) + " "
            embedd_each(row_text, api_key)
    return pd.read_excel(io.BytesIO(excel_content)).to_string()


def custom_ppt_reader(file, ppt_file_extension, ppt_file_content, api_key, rag):
    ppt_id = hash(file.filename)
    with open(f"{ppt_id}.{ppt_file_extension}", "wb") as outFile:
        outFile.write(ppt_file_content)
    pres = Presentation(f"{ppt_id}.{ppt_file_extension}")
    content = ""
    for slide in pres.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if len(shape.text) > 0:
                    content += f"{shape.text.strip()}\n"
                    slide_text += f"{shape.text.strip()}\n"
        if rag == "y":
            embedd_each(slide_text, api_key)
    os.remove(f"{ppt_id}.{ppt_file_extension}")
    return content


def history_maintenance(history, api_key):
    tokens = check_token_count(history)
    print("History Token Count: ")
    print(tokens)
    if tokens > 1000:
        prompt = (
            f"Instructions: Respond only with the same JSON formatting of the history. rephrase the content of each assistant response more"
            f"concise and"
            f"compressed. History: {history}")
        return customized_response(prompt, "[]", api_key)
    return None


def check_token_count(history):
    encoding = tiktoken.encoding_for_model("gpt-4-0125-preview")
    num_tokens = len(encoding.encode(history))
    return num_tokens
